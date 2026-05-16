"""
Generates the PowerPoint presentation for:
AI-Powered Handwriting Generation System
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ───────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD    = RGBColor(0x16, 0x2A, 0x3F)   # card navy
ACCENT     = RGBColor(0x00, 0xC8, 0xFF)   # electric blue
ACCENT2    = RGBColor(0x7B, 0x2F, 0xFF)   # violet
GREEN      = RGBColor(0x00, 0xE6, 0x76)   # success green
AMBER      = RGBColor(0xFF, 0xC1, 0x07)   # warning amber
RED        = RGBColor(0xFF, 0x45, 0x45)   # error red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xC5)
MID_GREY   = RGBColor(0x37, 0x47, 0x4F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # truly blank


# ── Helper utilities ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, color=WHITE, line_spacing=1.15):
    """lines = list of (text, bold, color_override_or_None)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        elif len(item) == 2:
            text, bold = item; col = color
        else:
            text, bold, col = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = col
    return txb


def bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)


def accent_bar(slide, y=Inches(0.55), h=Pt(3)):
    add_rect(slide, Inches(0.6), y, Inches(12.13), h, ACCENT)


def slide_title(slide, title, subtitle=None):
    bg(slide)
    accent_bar(slide)
    add_text(slide, title,
             Inches(0.6), Inches(0.65), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
                 font_size=16, color=ACCENT)


def card(slide, x, y, w, h):
    return add_rect(slide, x, y, w, h, BG_CARD)


def badge(slide, text, x, y, w=Inches(1.5), h=Inches(0.38), color=ACCENT):
    r = add_rect(slide, x, y, w, h, color)
    add_text(slide, text, x, y, w, h,
             font_size=11, bold=True, color=BG_DARK,
             align=PP_ALIGN.CENTER)


def metric_box(slide, label, value, unit, x, y, w=Inches(2.6), h=Inches(1.1),
               val_color=GREEN):
    card(slide, x, y, w, h)
    add_text(slide, label, x + Inches(0.12), y + Inches(0.05),
             w - Inches(0.2), Inches(0.3),
             font_size=11, color=LIGHT_GREY)
    add_text(slide, value, x + Inches(0.12), y + Inches(0.3),
             w - Inches(0.2), Inches(0.5),
             font_size=28, bold=True, color=val_color)
    add_text(slide, unit, x + Inches(0.12), y + Inches(0.78),
             w - Inches(0.2), Inches(0.25),
             font_size=10, color=LIGHT_GREY)


# ── TABLE helper ─────────────────────────────────────────────────────────────
def add_table(slide, headers, rows, x, y, w, h,
              col_widths=None, header_bg=ACCENT, header_fg=BG_DARK,
              font_size=13, header_font_size=13):
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    # column widths
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw
    # header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = Pt(header_font_size)
        run.font.color.rgb = header_fg
    # data rows
    for ri, row in enumerate(rows):
        row_bg = BG_CARD if ri % 2 == 0 else BG_DARK
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            if isinstance(val, tuple):
                run.text = val[0]
                run.font.color.rgb = val[1]
            else:
                run.text = str(val)
                run.font.color.rgb = WHITE
            run.font.size = Pt(font_size)
    return tbl


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
# big accent gradient stripe
add_rect(s, 0, Inches(2.5), SLIDE_W, Inches(2.8), BG_CARD)
add_rect(s, 0, Inches(2.5), Inches(0.18), Inches(2.8), ACCENT)

add_text(s, "AI-Powered Handwriting",
         Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Generation System",
         Inches(0.6), Inches(2.05), Inches(12), Inches(0.8),
         font_size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "Exploring 5 deep-learning approaches to generate realistic handwriting",
         Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.5),
         font_size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(s, "Abdullah Shahzad  ·  BSAI24082  ·  Programming for AI  ·  Semester 4",
         Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
         font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# decorative dots
for xi, col in [(Inches(5.5), ACCENT), (Inches(6.5), ACCENT2), (Inches(7.5), GREEN)]:
    add_rect(s, xi, Inches(6.0), Inches(0.15), Inches(0.15), col)


# =============================================================================
# SLIDE 2 — PROJECT OVERVIEW
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Project Overview", "What we built and why")
bg(s)
accent_bar(s)
add_text(s, "Project Overview",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "What we built and why",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

boxes = [
    (ACCENT,  "Goal",
     "Build a system that generates realistic handwriting for any letter (a–z, A–Z) and digit (0–9) in the style of real writers."),
    (ACCENT2, "Challenge",
     "Handwriting is highly personal — same character looks completely different across writers. Models must capture style, not just shape."),
    (GREEN,   "Approach",
     "We designed and compared 5 different deep-learning methods ranging from GANs to Bézier-curve regressors, finding the best trade-off between complexity and quality."),
]
for i, (col, title, body) in enumerate(boxes):
    bx = Inches(0.5) + i * Inches(4.2)
    card(s, bx, Inches(1.9), Inches(4.0), Inches(4.8))
    add_rect(s, bx, Inches(1.9), Inches(4.0), Inches(0.08), col)
    add_text(s, title, bx + Inches(0.2), Inches(2.05),
             Inches(3.6), Inches(0.4), font_size=16, bold=True, color=col)
    add_text(s, body, bx + Inches(0.2), Inches(2.5),
             Inches(3.6), Inches(3.8), font_size=13, color=LIGHT_GREY)

add_text(s, "62 unique characters  ·  6 real writers  ·  ~2,100 handwriting samples",
         Inches(0.6), Inches(7.0), Inches(12), Inches(0.35),
         font_size=12, color=MID_GREY, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 3 — DATASET
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Dataset", Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Real handwriting samples collected from 6 writers",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# stat boxes
stats = [
    ("6", "Writers"),
    ("~2,100", "PNG Samples"),
    ("62", "Characters"),
    ("~33", "Samples / char / writer"),
]
for i, (val, lbl) in enumerate(stats):
    bx = Inches(0.5) + i * Inches(3.2)
    metric_box(s, lbl, val, "", bx, Inches(1.9), w=Inches(2.9), h=Inches(1.2), val_color=ACCENT)

# writers list
card(s, Inches(0.5), Inches(3.3), Inches(5.8), Inches(3.8))
add_text(s, "Writers", Inches(0.7), Inches(3.45), Inches(5.4), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)
writers = [
    "writer_Abdullah",
    "writer_abdullah_60",
    "writer_Fatima",
    "writer_Hamza",
    "writer_Hashim",
    "writer_Salman_24067",
]
for i, w in enumerate(writers):
    add_text(s, f"  {w}", Inches(0.7), Inches(3.9) + i * Inches(0.5),
             Inches(5.3), Inches(0.45), font_size=13, color=WHITE)

# characters info
card(s, Inches(6.7), Inches(3.3), Inches(6.1), Inches(3.8))
add_text(s, "Character Set", Inches(6.9), Inches(3.45), Inches(5.7), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)
char_info = [
    ("Lowercase", "a – z", ACCENT),
    ("Uppercase", "A – Z", GREEN),
    ("Digits",    "0 – 9", AMBER),
]
for i, (cat, chars, col) in enumerate(char_info):
    y = Inches(3.95) + i * Inches(0.75)
    add_text(s, cat, Inches(6.9), y, Inches(2.2), Inches(0.4),
             font_size=13, color=LIGHT_GREY)
    add_text(s, chars, Inches(9.5), y, Inches(2.8), Inches(0.4),
             font_size=17, bold=True, color=col)

add_text(s, "File pattern:  WriterName_lc_a_r01.png  (lc = lowercase, uc = uppercase, n = digit)",
         Inches(6.9), Inches(6.0), Inches(5.7), Inches(0.8),
         font_size=11, color=LIGHT_GREY, italic=True)


# =============================================================================
# SLIDE 4 — 5 APPROACHES OVERVIEW
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "5 Approaches — Side by Side",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "We explored the solution space from pixel-level GANs to simple curve regressors",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

headers = ["#", "Name", "Generates?", "How", "Outcome"]
rows = [
    ["01", "GAN",
     ("Yes — pixels", WHITE),
     "Style encoder + Generator + AC-GAN Discriminator",
     ("Mode collapse", RED)],
    ["02", "MultiTaskCNN",
     ("No — classifier", AMBER),
     "Two heads: predicts char & writer from real images",
     ("57.6% char acc", AMBER)],
    ["03", "Font CNN",
     ("Yes — pixels", WHITE),
     "U-Net translates font glyph → handwriting style image",
     ("Blank output", RED)],
    ["04", "Bézier Reg",
     ("Yes — curves", WHITE),
     "CNN: reference photo → 6 Bézier curves → rendered",
     ("Mode collapse", RED)],
    ["05", "CharNet",
     ("Yes — curves", GREEN),
     "Char index → 6 Bézier curves → rendered  (no ref needed)",
     ("Best ✓", GREEN)],
]
cws = [Inches(0.45), Inches(1.6), Inches(1.4), Inches(4.9), Inches(1.9)]
add_table(s, headers, rows,
          Inches(0.5), Inches(1.85), Inches(12.3), Inches(5.3),
          col_widths=cws, font_size=12, header_font_size=13)


# =============================================================================
# SLIDE 5 — APPROACH 01: GAN (brief)
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Approach 01 — GAN",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Generative Adversarial Network  ·  pixel-level image synthesis",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# Architecture diagram (text-based)
components = [
    (Inches(0.5),  ACCENT,  "Style Encoder (CNN)", "Reads a real reference photo\nand compresses it into a\nstyle vector"),
    (Inches(4.3),  ACCENT2, "Generator",           "Takes style vector + char label\nand produces a full\nhandwriting image"),
    (Inches(8.1),  GREEN,   "Discriminator",       "AC-GAN: judges real vs fake\nand also predicts the\ncharacter label"),
]
for bx, col, title, desc in components:
    card(s, bx, Inches(1.95), Inches(3.5), Inches(2.8))
    add_rect(s, bx, Inches(1.95), Inches(3.5), Inches(0.07), col)
    add_text(s, title, bx + Inches(0.15), Inches(2.08),
             Inches(3.2), Inches(0.4), font_size=14, bold=True, color=col)
    add_text(s, desc, bx + Inches(0.15), Inches(2.55),
             Inches(3.2), Inches(2.0), font_size=12, color=LIGHT_GREY)

# arrows between boxes (simple text arrows)
for ax in [Inches(4.0), Inches(7.8)]:
    add_text(s, "→", ax, Inches(2.85), Inches(0.35), Inches(0.5),
             font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# results
card(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.1))
add_text(s, "Result", Inches(0.7), Inches(5.2), Inches(3), Inches(0.4),
         font_size=14, bold=True, color=RED)
result_lines = [
    ("Mode collapse: ", True, AMBER),
    ("the generator ignored the character label and always produced the same blurry blob regardless of input. "
     "Root cause — reference image was the same (dataset[0]) for all 62 characters; "
     "fixed in generate() but image quality remained poor.", False, LIGHT_GREY),
]
add_multiline(s, result_lines, Inches(0.7), Inches(5.6), Inches(11.8), Inches(1.4),
              font_size=12)


# =============================================================================
# SLIDE 6 — APPROACH 02: MultiTaskCNN
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Approach 02 — MultiTaskCNN",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Classification model — not a generator",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=AMBER)

# Left: what it does
card(s, Inches(0.5), Inches(1.85), Inches(5.5), Inches(5.3))
add_text(s, "What It Does", Inches(0.7), Inches(1.98), Inches(5), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)
points = [
    "• Takes a real handwriting image as input",
    "• Two classification heads:",
    "    – Character head: which of 62 characters?",
    "    – Writer head: which of 6 writers?",
    "",
    "• This is a classifier, not a generator.",
    "• Output = your own dataset photos with labels.",
    "• Useful for style-feature extraction and",
    "  verifying the dataset is well-structured.",
]
for i, p in enumerate(points):
    add_text(s, p, Inches(0.7), Inches(2.45) + i * Inches(0.48),
             Inches(5.1), Inches(0.45), font_size=12, color=LIGHT_GREY)

# Right: results
card(s, Inches(6.4), Inches(1.85), Inches(6.4), Inches(5.3))
add_text(s, "Results", Inches(6.6), Inches(1.98), Inches(5.8), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)

metrics = [
    ("Char Accuracy (Top-1)", "57.6 %", AMBER, "mediocre"),
    ("Char Accuracy (Top-5)", "89.9 %", GREEN, "good — right answer in top 5"),
    ("Writer Accuracy",       "31.9 %", RED,   "poor  (random chance = 16.7 %)"),
]
for i, (lbl, val, col, note) in enumerate(metrics):
    y = Inches(2.55) + i * Inches(1.1)
    add_text(s, lbl, Inches(6.6), y, Inches(3.5), Inches(0.4),
             font_size=11, color=LIGHT_GREY)
    add_text(s, val, Inches(6.6), y + Inches(0.35), Inches(2.5), Inches(0.5),
             font_size=22, bold=True, color=col)
    add_text(s, note, Inches(9.1), y + Inches(0.38), Inches(3.5), Inches(0.4),
             font_size=10, color=LIGHT_GREY, italic=True)

add_text(s, "⚠  Digits (0-8) all 0% F1 — model didn't learn digit classes",
         Inches(6.6), Inches(5.55), Inches(5.9), Inches(0.5),
         font_size=12, color=AMBER, bold=True)


# =============================================================================
# SLIDE 7 — APPROACH 03: Font CNN
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Approach 03 — Font to Handwriting CNN",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "U-Net translates a font-rendered glyph into a writer's style",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# Architecture
steps = [
    (ACCENT,  "Font Renderer",    "PIL renders the target\ncharacter in a clean\ncomputer font"),
    (ACCENT2, "U-Net Encoder",    "Downsamples the font\nglyph extracting spatial\nfeatures"),
    (GREEN,   "Style Injection",  "Writer style vector\ninjected at the bottleneck\nvia FiLM layers"),
    (AMBER,   "U-Net Decoder",    "Upsamples back to\nfull resolution with\nskip connections"),
]
for i, (col, title, desc) in enumerate(steps):
    bx = Inches(0.4) + i * Inches(3.22)
    card(s, bx, Inches(1.95), Inches(3.0), Inches(2.7))
    add_rect(s, bx, Inches(1.95), Inches(3.0), Inches(0.07), col)
    add_text(s, title, bx + Inches(0.12), Inches(2.07), Inches(2.8), Inches(0.4),
             font_size=13, bold=True, color=col)
    add_text(s, desc, bx + Inches(0.12), Inches(2.52), Inches(2.8), Inches(2.0),
             font_size=11, color=LIGHT_GREY)

for ax in [Inches(3.4), Inches(6.62), Inches(9.84)]:
    add_text(s, "→", ax, Inches(2.85), Inches(0.35), Inches(0.5),
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Results
card(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.35))
add_text(s, "Result", Inches(0.7), Inches(4.95), Inches(2), Inches(0.4),
         font_size=14, bold=True, color=RED)

res_cols = [
    (ACCENT,  "SSIM: 0.906", "Structurally matches\nreference (looks great\non paper)"),
    (RED,     "Ink coverage: 0.14%", "Nearly blank output —\nmodel generates almost\nno visible ink"),
    (RED,     "Distinctiveness: 0.0", "All 62 outputs identical —\ncomplete failure\nto differentiate characters"),
]
for i, (col, val, note) in enumerate(res_cols):
    bx = Inches(0.7) + i * Inches(4.15)
    add_text(s, val, bx, Inches(5.3), Inches(3.9), Inches(0.45),
             font_size=14, bold=True, color=col)
    add_text(s, note, bx, Inches(5.75), Inches(3.9), Inches(1.2),
             font_size=11, color=LIGHT_GREY)


# =============================================================================
# SLIDE 8 — APPROACH 04: Bézier Regression
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Approach 04 — Bézier Curve Regression",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "CNN predicts stroke paths instead of pixel values",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# Left column: how it works
card(s, Inches(0.5), Inches(1.85), Inches(6.1), Inches(5.3))
add_text(s, "How It Works", Inches(0.7), Inches(1.98), Inches(5.5), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)
steps_text = [
    ("1.  Reference photo of a character is fed in", False, LIGHT_GREY),
    ("2.  CNN backbone (ResNet-style) extracts features", False, LIGHT_GREY),
    ("3.  Regression head outputs 6 × 8 = 48 values", False, LIGHT_GREY),
    ("    (6 cubic Bézier curves × 4 control points × 2 coords)", False, LIGHT_GREY),
    ("4.  Curves are rendered onto a blank canvas", False, LIGHT_GREY),
    ("5.  Ground truth: real images → skeletonise → fit curves", False, LIGHT_GREY),
    ("6.  Loss = MSE between predicted and real curves", False, LIGHT_GREY),
]
add_multiline(s, steps_text, Inches(0.7), Inches(2.45), Inches(5.6), Inches(4.5),
              font_size=12)

# Right column: results
card(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(5.3))
add_text(s, "Results", Inches(7.2), Inches(1.98), Inches(5.2), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)

r_data = [
    ("SSIM", "0.720", GREEN, "reasonable structural similarity"),
    ("Ink Coverage", "0.55%", RED, "very low — sparse strokes"),
    ("Distinctiveness", "0.004", RED, "near-zero — mode collapse"),
]
for i, (lbl, val, col, note) in enumerate(r_data):
    y = Inches(2.55) + i * Inches(1.1)
    add_text(s, lbl, Inches(7.2), y, Inches(5.2), Inches(0.35),
             font_size=11, color=LIGHT_GREY)
    add_text(s, val, Inches(7.2), y + Inches(0.32), Inches(2.2), Inches(0.5),
             font_size=22, bold=True, color=col)
    add_text(s, note, Inches(9.4), y + Inches(0.35), Inches(3.2), Inches(0.45),
             font_size=10, color=LIGHT_GREY, italic=True)

add_text(s, "Issue: model always predicts the same average curve\nregardless of which character was requested.",
         Inches(7.2), Inches(5.5), Inches(5.4), Inches(0.8),
         font_size=11, color=AMBER, italic=True)


# =============================================================================
# SLIDE 9 — APPROACH 05: CharNet (architecture)
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Approach 05 — CharNet  (Best Model)",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Simplest approach — and the one that actually works",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=GREEN)

badge(s, "No reference image needed", Inches(8.5), Inches(1.2),
      w=Inches(4.3), color=GREEN)

# Pipeline
steps_c = [
    (ACCENT,  "Char Index\n(0 – 61)",  "Input: a single integer\nidentifying which character\nto generate"),
    (ACCENT2, "Embedding\nLayer",      "Maps char index to a\nlearned 64-dim vector\n(like word embeddings)"),
    (GREEN,   "MLP Head",              "3 fully-connected layers\npredict 48 floats\n(6 curves × 8 coords)"),
    (AMBER,   "Bézier\nRenderer",      "Control points drawn onto\na 128×128 canvas\nusing cubic splines"),
]
for i, (col, title, desc) in enumerate(steps_c):
    bx = Inches(0.4) + i * Inches(3.22)
    card(s, bx, Inches(1.95), Inches(3.0), Inches(2.9))
    add_rect(s, bx, Inches(1.95), Inches(3.0), Inches(0.08), col)
    add_text(s, title, bx + Inches(0.12), Inches(2.08), Inches(2.8), Inches(0.6),
             font_size=13, bold=True, color=col)
    add_text(s, desc, bx + Inches(0.12), Inches(2.72), Inches(2.8), Inches(1.9),
             font_size=11, color=LIGHT_GREY)

for ax in [Inches(3.4), Inches(6.62), Inches(9.84)]:
    add_text(s, "→", ax, Inches(2.9), Inches(0.35), Inches(0.5),
             font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Why it works better
card(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.1))
add_text(s, "Why It Works Better", Inches(0.7), Inches(5.2), Inches(5), Inches(0.4),
         font_size=13, bold=True, color=GREEN)
why = [
    ("• No style ambiguity: ", True, WHITE),
    ("each character index maps to one and only one output — no confusion from mismatched references.   ", False, LIGHT_GREY),
    ("• Simpler task: ", True, WHITE),
    ("the model only needs to learn 62 prototypical stroke shapes, not per-writer variation.   ", False, LIGHT_GREY),
    ("• Bézier representation: ", True, WHITE),
    ("predicting curves (not pixels) is a lower-dimensional, smoother regression target.", False, LIGHT_GREY),
]
add_multiline(s, why, Inches(0.7), Inches(5.65), Inches(11.8), Inches(1.3), font_size=12)


# =============================================================================
# SLIDE 10 — CharNet Results
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "CharNet — Evaluation Results",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Approach 05  ·  Trained for 105 epochs on T4 GPU",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# Top metric boxes
mets = [
    ("Overall MSE", "0.0042", "target < 0.005  ✓", GREEN),
    ("Endpoint Deviation", "0.074", "slightly above 0.05 target", AMBER),
    ("SSIM", "0.851", "strong structural similarity", GREEN),
    ("Sharpness", "0.731", "clean, crisp strokes", GREEN),
]
for i, (lbl, val, unit, col) in enumerate(mets):
    bx = Inches(0.4) + i * Inches(3.22)
    metric_box(s, lbl, val, unit, bx, Inches(1.85), w=Inches(3.0), h=Inches(1.2), val_color=col)

# Per-category breakdown
card(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.9))
add_text(s, "Per-Category Breakdown", Inches(0.7), Inches(3.42), Inches(6), Inches(0.38),
         font_size=13, bold=True, color=ACCENT)

cat_headers = ["Category", "Avg MSE", "Threshold", "Pass?", "Hardest Char", "Easiest Char"]
cat_rows = [
    [("Uppercase", AMBER),  "0.0046", "< 0.005", ("✓", GREEN), ("D  0.018", RED), ("L  0.0002", GREEN)],
    [("Lowercase", WHITE),  "0.0045", "< 0.005", ("✓", GREEN), ("e  0.017", RED), ("i  0.0003", GREEN)],
    [("Digits",    ACCENT), "0.0032", "< 0.005", ("✓", GREEN), ("9  0.011", AMBER), ("5  0.001", GREEN)],
]
add_table(s, cat_headers, cat_rows,
          Inches(0.5), Inches(3.82), Inches(12.3), Inches(3.2),
          col_widths=[Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.2), Inches(3.3), Inches(3.3)],
          font_size=12, header_font_size=12)


# =============================================================================
# SLIDE 11 — Image Quality Comparison
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Image Quality — All Generators",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Comparing the 4 approaches that actually generate new images",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

iq_headers = ["Approach", "SSIM ↑", "Ink Coverage", "Sharpness ↑", "Distinctiveness ↑", "Verdict"]
iq_rows = [
    [("01  GAN", WHITE),
     ("0.015", RED),
     ("84%  ← blurry blob", RED),
     ("0.72", GREEN),
     ("0.014", RED),
     ("Broken", RED)],
    [("03  Font CNN", WHITE),
     ("0.906", GREEN),
     ("0.14%  ← blank", RED),
     ("0.73", GREEN),
     ("0.000", RED),
     ("Blank output", RED)],
    [("04  Bézier Reg", WHITE),
     ("0.720", GREEN),
     ("0.55%  low", AMBER),
     ("0.64", AMBER),
     ("0.004", RED),
     ("Mode collapse", AMBER)],
    [("05  CharNet", GREEN),
     ("0.851", GREEN),
     ("2.5%  OK for curves", GREEN),
     ("0.73", GREEN),
     ("0.022", AMBER),
     ("Best ✓", GREEN)],
]
add_table(s, iq_headers, iq_rows,
          Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.5),
          col_widths=[Inches(1.7), Inches(1.3), Inches(2.5), Inches(1.5), Inches(2.3), Inches(3.0)],
          font_size=12, header_font_size=12)

add_text(s,
         "SSIM = Structural Similarity (1.0 = identical).  Ink Coverage = fraction of canvas covered by ink.  "
         "Distinctiveness = how different the 62 generated chars are from each other (0 = all identical).",
         Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
         font_size=10, color=LIGHT_GREY, italic=True)


# =============================================================================
# SLIDE 12 — Evaluation Framework
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Evaluation Framework",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "4 independent evaluation scripts covering every angle of model quality",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

scripts = [
    (ACCENT,  "eval_02_classifier.py",
     "Classification Report · Confusion Matrices · t-SNE by Writer",
     "Measures how well MultiTaskCNN identifies characters and writers"),
    (GREEN,   "eval_05_charnet.py",
     "Per-char MSE · Category MSE · Endpoint Deviation",
     "Compares predicted Bézier curves against skeletonised ground truth"),
    (ACCENT2, "eval_image_quality.py",
     "SSIM · Sharpness · Ink Coverage · Distinctiveness",
     "Pixel-level quality comparison across all 4 generator approaches"),
    (AMBER,   "eval_ocr_benchmark.py",
     "EasyOCR reads each generated char · Accuracy, F1 per approach",
     "Fairest cross-model benchmark — uses a neutral real-world OCR engine"),
]
for i, (col, name, outputs, desc) in enumerate(scripts):
    bx = (Inches(0.5) if i < 2 else Inches(6.7))
    by = Inches(1.95) + (i % 2) * Inches(2.6)
    card(s, bx, by, Inches(5.9), Inches(2.35))
    add_rect(s, bx, by, Inches(5.9), Inches(0.07), col)
    add_text(s, name, bx + Inches(0.15), by + Inches(0.12), Inches(5.5), Inches(0.4),
             font_size=13, bold=True, color=col)
    add_text(s, outputs, bx + Inches(0.15), by + Inches(0.55), Inches(5.5), Inches(0.5),
             font_size=11, color=WHITE)
    add_text(s, desc, bx + Inches(0.15), by + Inches(1.05), Inches(5.5), Inches(1.1),
             font_size=11, color=LIGHT_GREY, italic=True)

add_text(s, "All scripts write outputs to stats/outputs/ — charts + CSVs auto-generated.",
         Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
         font_size=11, color=MID_GREY, italic=True)


# =============================================================================
# SLIDE 13 — What Worked vs What Didn't
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "What Worked · What Didn't",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.65),
         font_size=32, bold=True)
add_text(s, "Honest assessment of each approach",
         Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT)

# Left: Worked
card(s, Inches(0.5), Inches(1.85), Inches(5.9), Inches(5.35))
add_rect(s, Inches(0.5), Inches(1.85), Inches(0.09), Inches(5.35), GREEN)
add_text(s, "What Worked", Inches(0.75), Inches(1.98), Inches(5.3), Inches(0.4),
         font_size=15, bold=True, color=GREEN)
worked = [
    "CharNet (05): MSE 0.0042 — below the 0.005 target",
    "CharNet SSIM 0.851 — strong structural quality",
    "CharNet sharpness 0.731 — clean crisp strokes",
    "MultiTaskCNN Top-5 accuracy 89.9% — right answer almost always in top 5",
    "Bézier approach: SSIM 0.72 — reasonable structural match",
    "Evaluation pipeline: 4 scripts, full metrics, reproducible",
]
for i, w in enumerate(worked):
    add_text(s, f"✓  {w}",
             Inches(0.75), Inches(2.5) + i * Inches(0.72),
             Inches(5.4), Inches(0.65), font_size=11.5, color=WHITE)

# Right: Didn't work
card(s, Inches(6.9), Inches(1.85), Inches(5.9), Inches(5.35))
add_rect(s, Inches(6.9), Inches(1.85), Inches(0.09), Inches(5.35), RED)
add_text(s, "What Didn't Work", Inches(7.15), Inches(1.98), Inches(5.3), Inches(0.4),
         font_size=15, bold=True, color=RED)
didnt = [
    "GAN: mode collapse — same blurry blob for all 62 chars",
    "Font CNN: generates nearly blank images (ink < 0.2%)",
    "Bézier Reg: mode collapse — same average curve every time",
    "MultiTaskCNN: digits 0-8 all 0% F1 — failed digit classes",
    "MultiTaskCNN writer accuracy 31.9% — barely above chance",
    "Distinctiveness across all models is low — chars too similar",
]
for i, d in enumerate(didnt):
    add_text(s, f"✗  {d}",
             Inches(7.15), Inches(2.5) + i * Inches(0.72),
             Inches(5.4), Inches(0.65), font_size=11.5, color=WHITE)


# =============================================================================
# SLIDE 14 — CONCLUSION
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, Inches(2.2), SLIDE_W, Inches(3.2), BG_CARD)
add_rect(s, 0, Inches(2.2), Inches(0.18), Inches(3.2), GREEN)
accent_bar(s, y=Inches(0.55))

add_text(s, "Conclusion",
         Inches(0.6), Inches(0.65), Inches(12), Inches(0.7),
         font_size=36, bold=True, color=WHITE)

add_text(s,
         "Simplicity beat complexity.",
         Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.65),
         font_size=26, bold=True, color=GREEN)

add_text(s,
         "CharNet — a single embedding layer + MLP — outperformed a GAN, a U-Net, and a CNN regression model "
         "on every meaningful metric. It achieves Bézier MSE of 0.0042 (target < 0.005), SSIM of 0.851, "
         "and crisp sharpness of 0.731 with no reference image required.",
         Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.3),
         font_size=14, color=LIGHT_GREY)

# Three takeaways
tks = [
    (ACCENT,  "Lesson 1 — Match the representation to the task",
     "Bézier curves are a natural fit for strokes. Pixels are not."),
    (ACCENT2, "Lesson 2 — Data quality beats model complexity",
     "All approaches shared the same 2,100-image dataset. Simpler models generalised better."),
    (GREEN,   "Lesson 3 — Evaluate honestly",
     "Mode collapse and blank outputs are real failures. Our eval scripts caught them objectively."),
]
for i, (col, title, body) in enumerate(tks):
    bx = Inches(0.5) + i * Inches(4.27)
    card(s, bx, Inches(5.05), Inches(4.0), Inches(2.15))
    add_rect(s, bx, Inches(5.05), Inches(4.0), Inches(0.07), col)
    add_text(s, title, bx + Inches(0.14), Inches(5.17), Inches(3.7), Inches(0.5),
             font_size=11, bold=True, color=col)
    add_text(s, body, bx + Inches(0.14), Inches(5.68), Inches(3.7), Inches(1.3),
             font_size=11, color=LIGHT_GREY)

add_text(s, "Abdullah Shahzad  ·  BSAI24082  ·  Programming for AI  ·  2026",
         Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
         font_size=11, color=MID_GREY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"D:\Semester_4\Programming for AI\Projects\AI-Powered-Handwriting-Generation-System\Handwriting_Generation_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
